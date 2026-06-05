"""
================================================================================
 Syntax Corporation © 2026 — All Rights Reserved
--------------------------------------------------------------------------------
 Project : EBS Contract Renewal PAF — BPA Renewal Automation Agent
 Module  : deliverables/build_sales_deck.py — exec/sales deck (.pptx)
 Version : 1.0.0      Build : 2026.06.04      Date : 2026-06-04
--------------------------------------------------------------------------------
 Generates a Syntax-branded exec deck into dist/, embedding the shared hero
 graphics and the one pricing model. python-pptx (no LibreOffice needed).
================================================================================
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

import pricing

ROOT = Path(__file__).resolve().parents[1]
DIST = ROOT / "deliverables" / "dist"
ASSETS = ROOT / "deliverables" / "assets"
DIST.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x06, 0x32, 0xA0); NAVY_DK = RGBColor(0x04, 0x1F, 0x66)
GREEN = RGBColor(0x3C, 0xC8, 0x5A); CYAN = RGBColor(0x1E, 0xB4, 0xE6)
GOLD = RGBColor(0xF1, 0xD4, 0x88); INK = RGBColor(0x16, 0x20, 0x3A)
SLATE = RGBColor(0x5B, 0x65, 0x77); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MIST = RGBColor(0xF4, 0xF7, 0xFC); LINE = RGBColor(0xDC, 0xE4, 0xF4)
HEAD = "Georgia"; BODY = "Calibri"

W, H = Inches(13.333), Inches(7.5)


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def _text(slide, x, y, w, h, text, size, color, *, bold=False, font=BODY,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return tb


def _tri(slide, x, y, s):
    # tri-color triangle motif (echoes the Syntax logo mark)
    t = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                               Inches(x), Inches(y), Inches(s), Inches(s))
    t.fill.solid(); t.fill.fore_color.rgb = NAVY; t.line.fill.background()
    t.shadow.inherit = False


def _footer(slide, n):
    _text(slide, Inches(0.5), H - Inches(0.42), Inches(7), Inches(0.3),
          "Syntax Corporation  ·  Confidential", 9, SLATE)
    _text(slide, W - Inches(1.0), H - Inches(0.42), Inches(0.5), Inches(0.3),
          str(n), 9, SLATE, align=PP_ALIGN.RIGHT)


def _kicker(slide, text):
    _tri(slide, 0.55, 0.5, 0.22)
    _text(slide, Inches(0.95), Inches(0.45), Inches(10), Inches(0.4),
          text.upper(), 12, NAVY, bold=True)


def _title(slide, text):
    _text(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.9),
          text, 30, NAVY_DK, bold=True, font=HEAD)


def content(prs, kicker, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE); _kicker(s, kicker); _title(s, title)
    return s


def card(slide, x, y, w, h, head, body, accent=NAVY):
    _rect(slide, x, y, w, h, MIST)
    _rect(slide, x, y, w, Inches(0.08), accent)
    _text(slide, x + Inches(0.18), y + Inches(0.18), w - Inches(0.36), Inches(0.5),
          head, 14, NAVY, bold=True, font=HEAD)
    _text(slide, x + Inches(0.18), y + Inches(0.72), w - Inches(0.36), h - Inches(0.9),
          body, 11.5, INK)


def build():
    m = pricing.model(); r = m["roi"]; impl = m["implementation"]
    msr = m["managed_services"]["terms"]
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H

    # --- 1. Title ---------------------------------------------------------
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, NAVY_DK)
    _rect(s, 0, H - Inches(0.6), W, Inches(0.6), NAVY)
    if (ASSETS / "syntax-logo.png").exists():
        s.shapes.add_picture(str(ASSETS / "syntax-logo.png"), Inches(0.6), Inches(0.55),
                             width=Inches(2.1))
    _text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(1.6),
          "Automating EBS Contract Renewals", 40, WHITE, bold=True, font=HEAD)
    _text(s, Inches(0.62), Inches(3.9), Inches(12), Inches(0.8),
          "An AI agent for Blanket Purchase Agreement renewals — built with the "
          "Oracle Private Agent Factory", 17, CYAN)
    _text(s, Inches(0.62), Inches(5.0), Inches(12), Inches(0.5),
          "Syntax Corporation  ·  Private Agent Factory series", 13, MIST)

    # --- 2. The problem ---------------------------------------------------
    s = content(prs, "The problem", "Renewals are a slow, error-prone bottleneck")
    cards = [("Manual & tedious", "Buyers re-key supplier price lists into EBS interface tables, "
              "agreement by agreement.", SLATE),
             ("Escalation creep", "Above-market price increases slip through; no consistent "
              "tolerance is applied.", GOLD),
             ("Coverage gaps", "Agreements lapse or overlap; renewals miss the contiguous term.", CYAN),
             ("No audit trail", "Hard to show why a renewal was accepted or held.", GREEN)]
    x = Inches(0.5)
    for head, body, acc in cards:
        card(s, x, Inches(2.0), Inches(2.95), Inches(2.6), head, body, acc)
        x += Inches(3.07)
    _footer(s, 2)

    # --- 3. The solution / pipeline --------------------------------------
    s = content(prs, "The solution", "One agent, end-to-end — read-only by default")
    if (DIST / "process.png").exists():
        s.shapes.add_picture(str(DIST / "process.png"), Inches(0.5), Inches(2.0),
                             width=Inches(12.3))
    _text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.2),
          "Ingest a supplier renewal quote → pull current prices from EBS via a governed MCP → "
          "apply deterministic renewal policy → stage a balanced PDOI batch for Import Price "
          "Catalogs. The buyer approves the renewed agreement in EBS — the agent never auto-approves "
          "a contract.", 13, INK)
    _footer(s, 3)

    # --- 4. Architecture --------------------------------------------------
    s = content(prs, "Architecture", "EBS 19c · PAF sidecar · managed MCP")
    if (DIST / "architecture.png").exists():
        s.shapes.add_picture(str(DIST / "architecture.png"), Inches(0.7), Inches(1.9),
                             width=Inches(11.9))
    _footer(s, 4)

    # --- 5. Policy --------------------------------------------------------
    s = content(prs, "Deterministic policy", "Auditable rules — outside the LLM")
    rows = [("Price escalation", "Hold only when the increase breaches BOTH a 7% cap AND a $250 "
             "annual-impact floor — real over-charges caught, trivial ones ignored."),
            ("Term", "Contiguous with the expiring agreement; length ≤ 36 months."),
            ("Validity", "Item still active/orderable; supplier site active."),
            ("Balance & QA", "Header AMOUNT_AGREED = Σ(price×qty); PASS / HOLD / FAIL gate before any load.")]
    y = Inches(2.0)
    for head, body in rows:
        _rect(s, Inches(0.5), y, Inches(0.12), Inches(0.95), NAVY)
        _text(s, Inches(0.8), y, Inches(11.8), Inches(0.4), head, 15, NAVY, bold=True, font=HEAD)
        _text(s, Inches(0.8), y + Inches(0.38), Inches(11.8), Inches(0.6), body, 12.5, INK)
        y += Inches(1.15)
    _footer(s, 5)

    # --- 6. Proof ---------------------------------------------------------
    s = content(prs, "Proof", "Verified live against EBS Vision")
    stats = [("51", "tests green\n(43 hermetic + 8 live)"),
             ("BPA 4467", "golden record\nHVAC Express, 5 lines"),
             ("$17,920", "balanced renewal\nbatch (Σ price×qty)"),
             ("100%", "SQL + write contract\nverified on the real DB")]
    x = Inches(0.5)
    for v, l in stats:
        _rect(s, x, Inches(2.2), Inches(2.95), Inches(2.3), MIST)
        _text(s, x, Inches(2.5), Inches(2.95), Inches(0.9), v, 30, NAVY, bold=True,
              font=HEAD, align=PP_ALIGN.CENTER)
        _text(s, x, Inches(3.5), Inches(2.95), Inches(0.9), l, 12, SLATE,
              align=PP_ALIGN.CENTER)
        x += Inches(3.07)
    _text(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.2),
          "Python↔PL/SQL parity validated live; PDOI batch round-tripped through the real "
          "PO_HEADERS_INTERFACE / PO_LINES_INTERFACE and rolled back clean.", 13, INK)
    _footer(s, 6)

    # --- 7. ROI -----------------------------------------------------------
    s = content(prs, "Business case", "Payback in weeks, not years")
    if (DIST / "roi_waterfall.png").exists():
        s.shapes.add_picture(str(DIST / "roi_waterfall.png"), Inches(0.5), Inches(2.0),
                             height=Inches(3.9))
    bx = Inches(7.2)
    for i, (v, l) in enumerate([(f"${r['annual_savings']:,.0f}", "net annual savings"),
                                (f"{r['payback_months']:.1f} mo", "payback"),
                                (f"{r['three_year_roi_pct']:.0f}%", "3-year ROI")]):
        yy = Inches(2.1) + Inches(1.35) * i
        _rect(s, bx, yy, Inches(5.4), Inches(1.15), MIST)
        _text(s, bx + Inches(0.2), yy + Inches(0.1), Inches(5), Inches(0.9),
              v, 26, GREEN if i != 1 else NAVY, bold=True, font=HEAD)
        _text(s, bx + Inches(2.4), yy + Inches(0.32), Inches(2.9), Inches(0.6),
              l, 14, SLATE)
    _footer(s, 7)

    # --- 8. Commercials ---------------------------------------------------
    s = content(prs, "Commercials", "Transparent, fixed-fee + managed services")
    _text(s, Inches(0.5), Inches(2.0), Inches(12), Inches(0.5),
          f"Fixed-fee implementation: ${impl['fixed_fee']:,.0f}  ·  {impl['timeline_weeks']} weeks  "
          f"·  ${impl['blended_rate']:.0f}/hr blended", 15, NAVY, bold=True)
    if (DIST / "tco.png").exists():
        s.shapes.add_picture(str(DIST / "tco.png"), Inches(0.5), Inches(2.7),
                             height=Inches(3.6))
    headers = [("12-month", msr["12"]), ("24-month", msr["24"]), ("36-month", msr["36"])]
    y = Inches(2.9)
    _text(s, Inches(7.2), Inches(2.5), Inches(5.5), Inches(0.4),
          "Managed services / month", 13, SLATE, bold=True)
    for i, (lab, t) in enumerate(headers):
        yy = y + Inches(0.95) * i
        _rect(s, Inches(7.2), yy, Inches(5.4), Inches(0.8), MIST)
        _text(s, Inches(7.4), yy + Inches(0.12), Inches(3), Inches(0.6), lab, 15, NAVY,
              bold=True, font=HEAD)
        _text(s, Inches(10.2), yy + Inches(0.1), Inches(2.3), Inches(0.6),
              f"${t['monthly']:,.0f}/mo", 16, INK, bold=True, align=PP_ALIGN.RIGHT)
        _text(s, Inches(10.2), yy + Inches(0.48), Inches(2.3), Inches(0.3),
              f"−{t['discount_pct']}%", 10, GREEN, align=PP_ALIGN.RIGHT)
    _footer(s, 8)

    # --- 9. Close ---------------------------------------------------------
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, NAVY_DK)
    if (ASSETS / "syntax-logo.png").exists():
        s.shapes.add_picture(str(ASSETS / "syntax-logo.png"), Inches(0.6), Inches(0.55),
                             width=Inches(2.0))
    _text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(1.4),
          "Renew every agreement — accurately, on time, on policy.", 32, WHITE,
          bold=True, font=HEAD)
    _text(s, Inches(0.62), Inches(4.1), Inches(12), Inches(0.8),
          "Start with a read-only value assessment against your EBS data.", 16, CYAN)
    _text(s, Inches(0.62), Inches(5.2), Inches(12), Inches(0.5),
          "Syntax Corporation © 2026  ·  Confidential", 12, MIST)

    out = DIST / "EBS_Contract_Renewal_PAF_Sales_Deck.pptx"
    prs.save(out); print(f"  OK  {out.name} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
